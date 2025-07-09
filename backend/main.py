from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
import os
import uuid
import shutil
from pathlib import Path
import json
from typing import Dict, Optional
import tempfile
import moviepy.editor as mp
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches
import whisper
from dotenv import load_dotenv

# Load environment variables from .env file
load_dotenv()

app = FastAPI(title="Voice-to-Slide Generator", version="1.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# In-memory job tracking
jobs: Dict[str, dict] = {}

# Load Whisper model
whisper_model = whisper.load_model("base")

# Initialize OpenAI client
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

# Create temp directories
UPLOAD_DIR = Path("backend/temp/uploads")
OUTPUT_DIR = Path("backend/temp/outputs")
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    if not file.filename.endswith('.mp4'):
        raise HTTPException(status_code=400, detail="Only MP4 files are supported")
    
    job_id = str(uuid.uuid4())
    file_path = UPLOAD_DIR / f"{job_id}.mp4"
    
    with open(file_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    jobs[job_id] = {
        "status": "uploaded",
        "filename": file.filename,
        "file_path": str(file_path),
        "progress": 10
    }
    
    print(f"File uploaded: {file.filename} -> {job_id}")
    return {"job_id": job_id, "message": "File uploaded successfully"}

@app.get("/status/{job_id}")
async def get_status(job_id: str):
    if job_id not in jobs:
        raise HTTPException(status_code=404, detail="Job not found")
    
    return jobs[job_id]

@app.get("/transcript/{job_id}")
async def get_transcript(job_id: str):
    if job_id not in jobs:
        raise HTTPException(status_code=404, detail="Job not found")
    
    job = jobs[job_id]
    
    if job["status"] == "uploaded":
        # Extract audio
        jobs[job_id]["status"] = "extracting_audio"
        jobs[job_id]["progress"] = 20
        
        try:
            video_path = job["file_path"]
            audio_path = UPLOAD_DIR / f"{job_id}.wav"
            
            video = mp.VideoFileClip(video_path)
            video.audio.write_audiofile(str(audio_path), verbose=False, logger=None)
            video.close()
            
            jobs[job_id]["status"] = "transcribing"
            jobs[job_id]["progress"] = 40
            jobs[job_id]["audio_path"] = str(audio_path)
            
            # Transcribe audio
            result = whisper_model.transcribe(str(audio_path))
            transcript = result["text"]
            
            jobs[job_id]["status"] = "transcript_ready"
            jobs[job_id]["progress"] = 60
            jobs[job_id]["transcript"] = transcript
            
            print(f"Transcription completed for job {job_id}")
            
        except Exception as e:
            jobs[job_id]["status"] = "error"
            jobs[job_id]["error"] = str(e)
            print(f"Error processing {job_id}: {str(e)}")
            raise HTTPException(status_code=500, detail=str(e))
    
    if "transcript" in jobs[job_id]:
        return {"transcript": jobs[job_id]["transcript"]}
    else:
        return {"message": "Transcript not ready yet"}

@app.post("/generate-slides/{job_id}")
async def generate_slides(job_id: str):
    if job_id not in jobs:
        raise HTTPException(status_code=404, detail="Job not found")
    
    job = jobs[job_id]
    
    if "transcript" not in job:
        raise HTTPException(status_code=400, detail="Transcript not available")
    
    try:
        jobs[job_id]["status"] = "generating_slides"
        jobs[job_id]["progress"] = 70
        
        # Generate slide content using OpenAI
        transcript = job["transcript"]
        
        prompt = f"""
        Convert the following transcript into a well-structured PowerPoint presentation outline.
        Create 5-8 slides with clear titles and bullet points.
        Format as JSON with this structure:
        {{
            "title": "Presentation Title",
            "slides": [
                {{
                    "title": "Slide Title",
                    "content": ["Bullet point 1", "Bullet point 2", "Bullet point 3"]
                }}
            ]
        }}
        
        Transcript: {transcript}
        """
        
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=1500
        )
        
        slide_content = json.loads(response.choices[0].message.content)
        
        jobs[job_id]["status"] = "creating_powerpoint"
        jobs[job_id]["progress"] = 85
        jobs[job_id]["slide_content"] = slide_content
        
        # Create PowerPoint presentation
        prs = Presentation()
        
        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = slide_content["title"]
        title_slide.placeholders[1].text = "Generated from audio transcript"
        
        # Content slides
        for slide_data in slide_content["slides"]:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_data["title"]
            
            content_shape = slide.placeholders[1]
            content_shape.text = "\n".join([f"• {point}" for point in slide_data["content"]])
        
        # Save PowerPoint file
        ppt_path = OUTPUT_DIR / f"{job_id}.pptx"
        prs.save(str(ppt_path))
        
        jobs[job_id]["status"] = "completed"
        jobs[job_id]["progress"] = 100
        jobs[job_id]["ppt_path"] = str(ppt_path)
        
        print(f"PowerPoint generated for job {job_id}")
        
        return {"message": "Slides generated successfully", "slide_content": slide_content}
        
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        print(f"Error generating slides for {job_id}: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/download/{job_id}")
async def download_file(job_id: str):
    if job_id not in jobs:
        raise HTTPException(status_code=404, detail="Job not found")
    
    job = jobs[job_id]
    
    if job["status"] != "completed" or "ppt_path" not in job:
        raise HTTPException(status_code=400, detail="PowerPoint not ready")
    
    ppt_path = job["ppt_path"]
    
    if not os.path.exists(ppt_path):
        raise HTTPException(status_code=404, detail="File not found")
    
    return FileResponse(
        path=ppt_path,
        filename=f"{job['filename']}.pptx",
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

@app.get("/")
async def root():
    return {"message": "Voice-to-Slide Generator API"}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)